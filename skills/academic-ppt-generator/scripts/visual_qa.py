#!/usr/bin/env python3
"""
Visual QA验证脚本 — 检测PPTX文件中的遮挡和溢出问题
使用python-pptx分析每个shape的坐标，检测潜在排版问题
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("❌ python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)

# 16:9宽屏页面尺寸（英寸）
PAGE_W = 13.33
PAGE_H = 7.5


def check_pptx(pptx_path, overlap_threshold=0.3, overflow_tolerance=0.1):
    """
    检查PPTX文件中的排版问题

    Args:
        pptx_path: PPTX文件路径
        overlap_threshold: 报告重叠的最小面积（平方英寸），默认0.3
        overflow_tolerance: 页面溢出容忍度（英寸），默认0.1

    Returns:
        dict: 检查结果，包含每页的问题列表
    """
    prs = Presentation(pptx_path)
    results = {}

    for i, slide in enumerate(prs.slides, 1):
        issues = []
        rects = []

        for j, shape in enumerate(slide.shapes):
            try:
                l = shape.left / 914400    # EMU to inches
                t = shape.top / 914400
                w = shape.width / 914400
                h = shape.height / 914400
                r = l + w
                b = t + h
                shape_type = str(shape.shape_type)

                # 检查页面溢出
                if r > PAGE_W + overflow_tolerance:
                    issues.append({
                        'type': 'overflow_right',
                        'shape': j,
                        'shape_type': shape_type,
                        'right_edge': round(r, 2),
                        'overflow': round(r - PAGE_W, 2),
                        'severity': 'high' if r - PAGE_W > 0.5 else 'low'
                    })
                if b > PAGE_H + overflow_tolerance:
                    issues.append({
                        'type': 'overflow_bottom',
                        'shape': j,
                        'shape_type': shape_type,
                        'bottom_edge': round(b, 2),
                        'overflow': round(b - PAGE_H, 2),
                        'severity': 'high' if b - PAGE_H > 0.5 else 'low'
                    })

                # 检查负坐标
                if l < -overflow_tolerance:
                    issues.append({
                        'type': 'negative_left',
                        'shape': j,
                        'shape_type': shape_type,
                        'left': round(l, 2),
                        'severity': 'low'
                    })
                if t < -overflow_tolerance:
                    issues.append({
                        'type': 'negative_top',
                        'shape': j,
                        'shape_type': shape_type,
                        'top': round(t, 2),
                        'severity': 'low'
                    })

                rects.append((l, t, r, b, j, shape_type))
            except Exception:
                continue

        # 检查显著重叠
        for j_idx in range(len(rects)):
            for k_idx in range(j_idx + 1, len(rects)):
                l1, t1, r1, b1, j, j_type = rects[j_idx]
                l2, t2, r2, b2, k, k_type = rects[k_idx]

                ox = min(r1, r2) - max(l1, l2)
                oy = min(b1, b2) - max(t1, t2)

                if ox > 0.08 and oy > 0.08:
                    overlap_area = ox * oy
                    if overlap_area > overlap_threshold:
                        # 判断是否为设计意图内的重叠（如背景矩形）
                        is_background = (
                            'RECTANGLE' in j_type or 'RECTANGLE' in k_type
                        ) and overlap_area > 5.0

                        issues.append({
                            'type': 'overlap',
                            'shape_a': j,
                            'shape_b': k,
                            'overlap_area': round(overlap_area, 2),
                            'overlap_x': round(ox, 2),
                            'overlap_y': round(oy, 2),
                            'likely_intentional': is_background,
                            'severity': 'info' if is_background else 'medium'
                        })

        results[f'slide_{i}'] = {
            'total_shapes': len(rects),
            'issues': issues,
            'has_problems': any(
                iss.get('severity') in ['high', 'medium']
                for iss in issues
                if not iss.get('likely_intentional', False)
            )
        }

    return results


def print_report(results):
    """打印检查报告"""
    total_issues = 0
    problem_slides = 0

    print('=' * 60)
    print('📊 PPTX Visual QA Report')
    print('=' * 60)
    print(f'Page size: {PAGE_W}" × {PAGE_H}" (16:9)')
    print()

    for slide_key, data in results.items():
        slide_num = slide_key.replace('slide_', '')
        issues = data['issues']

        if not issues:
            print(f'✅ Slide {slide_num}: {data["total_shapes"]} shapes, no issues')
            continue

        # 过滤掉设计意图内的重叠
        real_issues = [iss for iss in issues if not iss.get('likely_intentional', False)]

        if real_issues:
            problem_slides += 1
            total_issues += len(real_issues)
            print(f'⚠️  Slide {slide_num}: {data["total_shapes"]} shapes, {len(real_issues)} issues')
            for iss in real_issues:
                if iss['type'] == 'overflow_right':
                    print(f'    → OVERFLOW RIGHT: Shape {iss["shape"]} ({iss["shape_type"]}) '
                          f'right edge at {iss["right_edge"]}" (overflow {iss["overflow"]}")')
                elif iss['type'] == 'overflow_bottom':
                    print(f'    → OVERFLOW BOTTOM: Shape {iss["shape"]} ({iss["shape_type"]}) '
                          f'bottom edge at {iss["bottom_edge"]}" (overflow {iss["overflow"]}")')
                elif iss['type'] == 'overlap':
                    print(f'    → OVERLAP: Shape {iss["shape_a"]} & {iss["shape_b"]} '
                          f'area={iss["overlap_area"]} sq.in.')
        else:
            print(f'✅ Slide {slide_num}: {data["total_shapes"]} shapes, '
                  f'{len(issues)} intentional overlaps only')

    print()
    print('=' * 60)
    if total_issues == 0:
        print('✅ ALL CLEAR — No layout issues detected!')
    else:
        print(f'⚠️  Found {total_issues} issues in {problem_slides} slides')
    print('=' * 60)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python visual_qa.py <pptx_path> [--overlap-threshold 0.3]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    threshold = 0.3

    if '--overlap-threshold' in sys.argv:
        idx = sys.argv.index('--overlap-threshold')
        threshold = float(sys.argv[idx + 1])

    if not Path(pptx_path).exists():
        print(f"❌ File not found: {pptx_path}")
        sys.exit(1)

    results = check_pptx(pptx_path, overlap_threshold=threshold)
    print_report(results)

    # Export JSON for programmatic use
    json_path = Path(pptx_path).with_suffix('.qa.json')
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2, ensure_ascii=False)
    print(f'\n📄 Detailed report saved to: {json_path}')
